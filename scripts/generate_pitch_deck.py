from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "BetCode-Bridge-Company-Pitch.pptx"

BG = RGBColor(0x0B, 0x12, 0x20)
PANEL = RGBColor(0x11, 0x1B, 0x2E)
ACCENT = RGBColor(0xD4, 0xAF, 0x37)
TEXT = RGBColor(0xF4, 0xF7, 0xFF)
MUTED = RGBColor(0xB8, 0xC2, 0xD9)
GOOD = RGBColor(0x22, 0xC5, 0x5E)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title):
    shape = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.6))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TEXT


def add_footer(slide, text="BetCode Bridge - Confidential Partner Brief"):
    shape = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3))
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items, x=0.9, y=1.35, w=8.5, h=5.3, font_size=20):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(12)
        p.bullet = True


def add_panel(slide, x, y, w, h, color=PANEL):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0x2A, 0x3A, 0x52)
    return shape


def add_center_text(slide, text, x, y, w, h, size=18, color=TEXT, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_panel(s, 0.9, 0.9, 11.5, 5.6)
    add_center_text(s, "BETCODE BRIDGE", 1.2, 1.45, 10.8, 0.8, size=44, color=ACCENT, bold=True)
    add_center_text(s, "API Partnership Proposal for Bookmaker Growth", 1.2, 2.35, 10.8, 0.45, size=22)
    add_center_text(s, "Focus: Reactivating inactive customers and increasing repeat betting sessions", 1.2, 3.05, 10.8, 0.5, size=16, color=MUTED)
    add_center_text(s, f"Meeting Deck | {date.today().year}", 4.0, 4.45, 5.3, 0.5, size=16, bold=True)
    add_footer(s, "BetCode Bridge - Partner Meeting Deck")

    # Slide 2: Opportunity
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "1. The Opportunity for Your Bookmaker")
    add_bullets(
        s,
        [
            "Many customers become inactive because bet-slip recreation is slow and frustrating.",
            "BetCode Bridge removes friction by turning shared slips into actionable bets quickly.",
            "Lower friction drives return visits, stronger retention, and better conversion from dormant users.",
            "This is a practical reactivation channel without a major product overhaul on your side.",
        ],
        x=0.9,
        y=1.4,
        w=7.5,
        h=4.9,
        font_size=16,
    )
    add_panel(s, 8.95, 1.55, 3.3, 2.0)
    add_center_text(s, "Primary Benefit", 9.2, 1.78, 2.8, 0.3, size=13, color=ACCENT, bold=True)
    add_center_text(s, "Reawaken\ninactive\ncustomers", 9.2, 2.08, 2.8, 1.25, size=18, bold=True)
    add_panel(s, 8.95, 3.8, 3.3, 1.8)
    add_center_text(s, "Secondary Gains", 9.2, 4.02, 2.8, 0.3, size=13, color=ACCENT, bold=True)
    add_center_text(s, "More repeat sessions\nHigher usage", 9.2, 4.3, 2.8, 0.95, size=15, bold=True)
    add_footer(s)

    # Slide 3: How it works
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "2. How BetCode Bridge Works")

    x_positions = [0.85, 3.8, 6.75, 9.7]
    titles = ["User Inputs", "Decode", "Map & Compare", "Output"]
    bodies = [
        "Share code\nand source\nbookmaker",
        "Source slip\nnormalized\ninto selections",
        "Markets\nand odds\nmapped",
        "Destination\ncode and\ninsights",
    ]

    for i in range(4):
        add_panel(s, x_positions[i], 2.2, 2.35, 2.3)
        add_center_text(s, titles[i], x_positions[i] + 0.1, 2.45, 2.15, 0.35, size=14, color=ACCENT, bold=True)
        add_center_text(s, bodies[i], x_positions[i] + 0.15, 3.0, 2.05, 1.3, size=16)
        if i < 3:
            add_center_text(s, ">", x_positions[i] + 2.45, 3.15, 0.35, 0.55, size=28, color=ACCENT, bold=True)

    add_center_text(
        s,
        "Current build is integration-ready and can connect to your official API with your preferred auth model.",
        1.0,
        5.45,
        11.4,
        0.45,
        size=14,
        color=MUTED,
    )
    add_footer(s)

    # Slide 4: Commercial impact
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "3. Commercial Impact for Your Team")
    add_panel(s, 0.85, 1.35, 12.0, 1.1)
    add_center_text(
        s,
        "Core value: reawaken inactive users and turn slip intent into active betting behavior.",
        1.1,
        1.7,
        11.5,
        0.4,
        size=19,
        color=ACCENT,
        bold=True,
    )

    cards = [
        ("Reactivation", "Dormant users return\nwith lower friction", GOOD),
        ("Retention", "More repeat\nweekly sessions", ACCENT),
        ("Engagement", "Faster slip completion\nand more interactions", RGBColor(0xF5, 0x9E, 0x0B)),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 0.9 + (i * 4.1)
        add_panel(s, x, 2.9, 3.75, 3.1)
        bar = s.shapes.add_shape(1, Inches(x + 0.22), Inches(3.2), Inches(0.15), Inches(2.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        box = s.shapes.add_textbox(Inches(x + 0.52), Inches(3.2), Inches(3.0), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT
        box2 = s.shapes.add_textbox(Inches(x + 0.52), Inches(3.85), Inches(3.0), Inches(1.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(15)
        p2.font.color.rgb = MUTED

    add_footer(s)

    # Slide 5: API request
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "4. API Access Request")
    add_bullets(
        s,
        [
            "Sandbox and production base URLs",
            "Authentication method (Bearer token, API key, or signed headers)",
            "Share-code or bet-slip decode endpoint",
            "Request and response schema with sample payloads",
            "Rate limits, quotas, and retry behavior",
            "Onboarding and production approval checklist",
        ],
        x=0.9,
        y=1.45,
        w=7.8,
        h=4.9,
        font_size=16,
    )
    add_panel(s, 8.95, 1.55, 3.5, 4.65)
    add_center_text(s, "Reference Response", 9.25, 1.85, 2.9, 0.35, size=13, color=ACCENT, bold=True)
    code = (
        '{\n'
        '  "code": "0",\n'
        '  "message": "success",\n'
        '  "data": {\n'
        '    "shareCode": "waxec6",\n'
        '    "outcomes": [...]\n'
        '  }\n'
        '}'
    )
    box = s.shapes.add_textbox(Inches(9.2), Inches(2.25), Inches(3.0), Inches(3.5))
    p = box.text_frame.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT
    add_footer(s)

    # Slide 6: Security
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "5. Security and Compliance Approach")
    add_bullets(
        s,
        [
            "Credentials are stored in environment variables, never hardcoded.",
            "Access can be sandbox-first with controlled rollout to production.",
            "Integration can be toggled off safely if the API is unavailable.",
            "Usage logs can be limited to non-sensitive metadata.",
            "Architecture supports least-privilege access and scoped tokens.",
        ],
        x=0.9,
        y=1.45,
        w=8.0,
        h=4.9,
        font_size=16,
    )
    add_panel(s, 9.15, 1.75, 3.35, 3.9)
    add_center_text(s, "Risk-Control\nHighlights", 9.45, 2.25, 2.75, 0.9, size=18, color=GOOD, bold=True)
    add_center_text(s, "No credential leak\nSafe fallback mode\nPredictable integration", 9.45, 3.2, 2.75, 1.6, size=14, color=MUTED)
    add_footer(s)

    # Slide 7: Pilot plan
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "6. Pilot Plan and Next Steps")
    t = s.shapes.add_textbox(Inches(0.95), Inches(1.25), Inches(6.0), Inches(0.5))
    p = t.text_frame.paragraphs[0]
    p.text = "Proposed 3-phase rollout"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    phases = [
        ("Phase 1", "Sandbox Integration", "Validate auth, endpoint contract, and response mapping."),
        ("Phase 2", "Controlled Pilot", "Small user cohort, monitor reliability and engagement lift."),
        ("Phase 3", "Production Scale", "Wider rollout with reporting and optimization loop."),
    ]
    for i, (phase, title, detail) in enumerate(phases):
        y = 2.0 + i * 1.62
        add_panel(s, 0.95, y, 11.6, 1.25)
        b1 = s.shapes.add_textbox(Inches(1.2), Inches(y + 0.32), Inches(1.2), Inches(0.35))
        p1 = b1.text_frame.paragraphs[0]
        p1.text = phase
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = GOOD

        b2 = s.shapes.add_textbox(Inches(2.45), Inches(y + 0.18), Inches(3.8), Inches(0.4))
        p2 = b2.text_frame.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = TEXT

        b3 = s.shapes.add_textbox(Inches(2.45), Inches(y + 0.62), Inches(9.8), Inches(0.42))
        p3 = b3.text_frame.paragraphs[0]
        p3.text = detail
        p3.font.size = Pt(14)
        p3.font.color.rgb = MUTED

    add_footer(s)

    # Slide 8: Close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_center_text(s, "Thank You", 0.6, 2.0, 12.1, 0.9, size=54, color=ACCENT, bold=True)
    add_center_text(
        s,
        "Let us activate dormant users and grow repeat betting sessions together.",
        1.0,
        3.05,
        11.3,
        0.55,
        size=21,
    )
    add_panel(s, 3.9, 4.25, 5.5, 1.6)
    add_center_text(s, "Contact", 4.15, 4.55, 5.0, 0.3, size=15, color=ACCENT, bold=True)
    add_center_text(s, "[Your Name]\n[Company/App Name]\n[Email Address]\n[Phone Number]", 4.15, 4.92, 5.0, 0.88, size=13)
    add_footer(s, "BetCode Bridge - Ready for Partner Pilot")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Presentation created: {OUTPUT}")


if __name__ == "__main__":
    main()
